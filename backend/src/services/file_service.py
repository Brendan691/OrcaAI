"""文件存储与解析服务

存储层做成可切换抽象(见 ADR-0002):
- 本地默认:LocalFileStorage —— 存到 data/files/,零外部依赖
- 生产可选:MinioFileStorage —— 延迟 import minio,仅在 STORAGE_BACKEND=minio 时加载

解析层:轻量库(pypdf / python-docx / python-pptx)直接解析,不引入重型的 markitdown。
"""
import io
import uuid
from pathlib import Path
from typing import Protocol

from ..core.config import settings


class FileStorage(Protocol):
    """文件存储接口。任何后端(本地 / MinIO / S3)都实现这几个方法。"""

    def save(self, data: bytes, filename: str, content_type: str = "") -> str:
        """保存文件,返回 object_name(用于后续读取/删除)。"""
        ...

    def load(self, object_name: str) -> bytes:
        ...

    def delete(self, object_name: str) -> None:
        ...

    def get_url(self, object_name: str, expires: int = 3600) -> str:
        """返回可访问文件的地址。"""
        ...


class LocalFileStorage:
    """本地文件夹存储(默认)。文件落在 settings.FILE_STORAGE_DIR。"""

    def __init__(self):
        self.root = Path(settings.FILE_STORAGE_DIR)
        self.root.mkdir(parents=True, exist_ok=True)

    def save(self, data: bytes, filename: str, content_type: str = "") -> str:
        ext = Path(filename).suffix.lower()
        object_name = f"{uuid.uuid4().hex}{ext}"
        (self.root / object_name).write_bytes(data)
        return object_name

    def load(self, object_name: str) -> bytes:
        return (self.root / object_name).read_bytes()

    def delete(self, object_name: str) -> None:
        target = self.root / object_name
        if target.exists():
            target.unlink()

    def get_url(self, object_name: str, expires: int = 3600) -> str:
        # 本地存储通过后端下载接口访问
        return f"/api/files/{object_name}/raw"


class MinioFileStorage:
    """MinIO / S3 存储(生产)。延迟 import,避免本地没装 minio 时报错。"""

    def __init__(self):
        from minio import Minio  # 延迟导入:仅生产用

        self.bucket = settings.S3_BUCKET
        self._client = Minio(
            settings.S3_ENDPOINT,
            access_key=settings.S3_ACCESS_KEY_ID,
            secret_key=settings.S3_SECRET_ACCESS_KEY,
            secure=settings.S3_SECURE,
        )
        if not self._client.bucket_exists(self.bucket):
            self._client.make_bucket(self.bucket)

    def save(self, data: bytes, filename: str, content_type: str = "") -> str:
        ext = Path(filename).suffix.lower()
        object_name = f"{uuid.uuid4().hex}{ext}"
        self._client.put_object(
            self.bucket, object_name, io.BytesIO(data), len(data),
            content_type=content_type or "application/octet-stream",
        )
        return object_name

    def load(self, object_name: str) -> bytes:
        resp = self._client.get_object(self.bucket, object_name)
        try:
            return resp.read()
        finally:
            resp.close()
            resp.release_conn()

    def delete(self, object_name: str) -> None:
        self._client.remove_object(self.bucket, object_name)

    def get_url(self, object_name: str, expires: int = 3600) -> str:
        from datetime import timedelta
        return self._client.presigned_get_object(
            self.bucket, object_name, expires=timedelta(seconds=expires)
        )


def _make_storage() -> FileStorage:
    """按配置选择存储后端。"""
    if settings.STORAGE_BACKEND == "minio":
        return MinioFileStorage()
    return LocalFileStorage()


class FileService:
    """文件服务:存储(委托给 FileStorage)+ 多格式解析为纯文本。"""

    SUPPORTED_EXTENSIONS = {
        ".pdf", ".doc", ".docx", ".ppt", ".pptx",
        ".txt", ".md", ".csv", ".json", ".xml",
    }

    def __init__(self):
        self._storage: FileStorage | None = None

    @property
    def storage(self) -> FileStorage:
        # 延迟创建:import 时不触碰存储后端
        if self._storage is None:
            self._storage = _make_storage()
        return self._storage

    def upload_file(self, file_data: bytes, filename: str, content_type: str = "") -> str:
        return self.storage.save(file_data, filename, content_type)

    def download_file(self, object_name: str) -> bytes:
        return self.storage.load(object_name)

    def delete_file(self, object_name: str) -> None:
        self.storage.delete(object_name)

    def get_file_url(self, object_name: str, expires: int = 3600) -> str:
        return self.storage.get_url(object_name, expires)

    def parse_to_text(self, file_data: bytes, filename: str) -> str:
        """把文件解析为纯文本,用于切片和向量化。无法解析时返回以 [无法解析 开头的提示。"""
        ext = Path(filename).suffix.lower()

        if ext in (".txt", ".md", ".csv", ".json", ".xml"):
            return file_data.decode("utf-8", errors="replace")

        if ext == ".pdf":
            try:
                from pypdf import PdfReader
                reader = PdfReader(io.BytesIO(file_data))
                return "\n".join(page.extract_text() or "" for page in reader.pages)
            except Exception as e:
                return f"[无法解析 PDF: {e}]"

        if ext in (".doc", ".docx"):
            try:
                from docx import Document as DocxDocument
                doc = DocxDocument(io.BytesIO(file_data))
                return "\n".join(p.text for p in doc.paragraphs)
            except Exception as e:
                return f"[无法解析 Word: {e}]"

        if ext in (".ppt", ".pptx"):
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_data))
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            texts.append(shape.text)
                return "\n".join(texts)
            except Exception as e:
                return f"[无法解析 PPT: {e}]"

        return f"[无法解析的文件类型: {ext}]"

    @staticmethod
    def is_supported(filename: str) -> bool:
        return Path(filename).suffix.lower() in FileService.SUPPORTED_EXTENSIONS


file_service = FileService()

